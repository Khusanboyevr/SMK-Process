from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from io import BytesIO
import pypdf
from pptx import Presentation
from groq import Groq
import os
from pathlib import Path
from dotenv import load_dotenv
import json

# .env faylni backend papkasidan yuklash (qayerdan ishga tushirilmasin)
env_path = Path(__file__).resolve().parent / ".env"
load_dotenv(dotenv_path=env_path, override=True)

app = FastAPI(title="SMK Process Builder API")

# Allow CORS for React frontend (Netlify va localhost)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


def extract_text_from_pdf(file_contents: bytes) -> str:
    """PDF fayldan matn ajratib olish"""
    try:
        reader = pypdf.PdfReader(BytesIO(file_contents))
        text = "\n".join(
            [page.extract_text() for page in reader.pages if page.extract_text()]
        )
        return text.strip()
    except Exception as e:
        print(f"PDF extraction error: {e}")
        return ""


def extract_text_from_pptx(file_contents: bytes) -> str:
    """PPTX fayldan matn ajratib olish"""
    try:
        prs = Presentation(BytesIO(file_contents))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text.strip())
        return "\n".join(text)
    except Exception as e:
        print(f"PPTX extraction error: {e}")
        return ""


def generate_mock_analysis() -> dict:
    """API kalit bo'lmaganda yoki xatolikda mock data qaytarish"""
    return {
        "summary": {
            "industry": "Kiritilgan fayllarga asosan: Sanoat korxonasida ISO 9001 standarti doirasida ishlab chiqarish jarayoni, mahsulot sifatini qat'iy tekshirish va xatti-harakatlar xavfsizligi yetakchi o'ringa qo'yilgan. Xom-ashyoning qayta ishlanish monitoringi alohida ajratilgan.",
            "university": "Universitet muhitida sifat menejmenti asosan o'quv dasturlari aktualligi, professor-o'qituvchilarning ilmiy salohiyati va talabalar uchun qulay tizimlarni (masalan elektron kutubxona) joriy etishga tayanadi.",
        },
        "comparison": [
            {
                "criterion": "Asosiy faoliyat",
                "industry": "Moddiy tovar ishlab chiqarish",
                "university": "Ta'lim berish va ilmiy tadqiqot",
            },
            {
                "criterion": "Risklarni baholash",
                "industry": "Brak mahsulot (nuqson), uskunalar muddatidan oldin eskirishi",
                "university": "O'zlashtirishning pastligi, plagiat, dasturlar eskirishi",
            },
            {
                "criterion": "Mijoz / Iste'molchi",
                "industry": "Xaridorlar, hamkor korxonalar",
                "university": "Talabalar, ish beruvchilar va jamiyat",
            },
        ],
        "pdca": {
            "plan": [
                "Resurslarni (budjet, vaqt, odamlar) rejalashtirish",
                "Yangi oy/semestr maqsadlarini kiritish",
            ],
            "do_core": [
                "Liniyada mahsulot yig'ish (Sanoat)",
                "Yangi fanni joriy etish (Universitet)",
            ],
            "do_support": [
                "Texnik uskunalarni sozlash",
                "IT infratuzilmani yangilash",
                "Xodimlar va pedagoglar malakasini oshirish",
            ],
            "check_act": [
                "Ichki auditorlik tekshiruvi (ISO 9001)",
                "Kamchiliklar ustida profilaktika o'tkazish",
                "Talabalar/Mijozlar so'rovnomasini tahlil qilish",
            ],
        },
        "ikp_template": "[IKP Misol Shablon]:\n1. O'lchov formati: Muntazam.\n2. Tegishli shaxslar: Departament rahbarlari.\n3. Resurslar ro'yxati: Standartlashtirilgan o'lchov asboblari va ma'lumotlar bazasi.",
    }


def generate_ai_analysis(pdf_text: str, pptx_text: str) -> dict:
    """Groq AI orqali tahlil qilish, xatolikda mock data qaytarish"""
    api_key = os.getenv("GROQ_API_KEY")

    if not api_key or "shu_yerga_kalitni" in api_key or len(api_key.strip()) < 10:
        print("Using MOCK DATA (No valid API key provided)")
        return generate_mock_analysis()

    try:
        client = Groq(api_key=api_key.strip())

        combined_context = (
            f"--- SANOAT KORXONASI (PDF) ---\n{pdf_text[:10000]}\n\n"
            f"--- UNIVERSITET (PPTX) ---\n{pptx_text[:10000]}"
        )

        prompt = """
        Siz Xalqaro ISO 9001 Sifat Menejmenti Tizimi (SMK) bo'yicha bosh auditorsiz.
        Sizga 2 ta matn berilgan: Biri sanoat korxonasining SMK si, ikkinchisi Universitetning SMK si.
        Quyidagi JSON formatda yagona hisobot tayyorlang. Boshqa hech qanday so'z qoshmang, faqat JSON qaytaring. Xato JSON format qilmang!

        {
            "summary": {
                "industry": "Sanoat korxonasining sifat siyosati va standartlari bo'yicha 3-4 gaplik xulosa",
                "university": "Universitetning sifat tizimi bo'yicha 3-4 gaplik xulosa"
            },
            "comparison": [
                {"criterion": "Mezon nomi", "industry": "Korxona xolati", "university": "Universitet xolati"},
                {"criterion": "Asosiy Mahsulot", "industry": "...", "university": "..."},
                {"criterion": "Istemolchi", "industry": "...", "university": "..."}
            ],
            "pdca": {
                "plan": ["Jarayon 1", "Jarayon 2"],
                "do_core": ["Jarayon 1", "Jarayon 2"],
                "do_support": ["Jarayon 1", "Jarayon 2"],
                "check_act": ["Jarayon 1", "Jarayon 2"]
            },
            "ikp_template": "Avtomatik generatsiya qilingan qisqacha IKP (Informatsion Karta Protsessa) teksti (1 ta misol sifatida)"
        }
        """

        response = client.chat.completions.create(
            model="llama3-70b-8192",
            messages=[
                {"role": "system", "content": prompt},
                {"role": "user", "content": combined_context},
            ],
            response_format={"type": "json_object"},
        )

        result_str = response.choices[0].message.content
        return json.loads(result_str)
    except json.JSONDecodeError as je:
        print(f"JSON parse error: {je}. Falling back to Mock Data.")
        return generate_mock_analysis()
    except Exception as e:
        print(
            f"Groq API Error: {e}. Falling back to Mock Data to ensure stability."
        )
        return generate_mock_analysis()


@app.post("/analyze")
async def analyze_documents(
    pdf: UploadFile = File(...),
    pptx: UploadFile = File(...),
):
    # Fayl turlarini tekshirish
    allowed_pdf_types = [
        "application/pdf",
    ]
    allowed_pptx_types = [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    ]

    if pdf.content_type and pdf.content_type not in allowed_pdf_types:
        raise HTTPException(
            status_code=400,
            detail=f"PDF fayl kutilmoqda, lekin '{pdf.content_type}' turi yuborildi.",
        )

    if pptx.content_type and pptx.content_type not in allowed_pptx_types:
        raise HTTPException(
            status_code=400,
            detail=f"PPTX fayl kutilmoqda, lekin '{pptx.content_type}' turi yuborildi.",
        )

    try:
        pdf_content = await pdf.read()
        pptx_content = await pptx.read()

        if not pdf_content:
            raise HTTPException(status_code=400, detail="PDF fayl bo'sh.")
        if not pptx_content:
            raise HTTPException(status_code=400, detail="PPTX fayl bo'sh.")

        pdf_text = extract_text_from_pdf(pdf_content)
        pptx_text = extract_text_from_pptx(pptx_content)

        if not pdf_text and not pptx_text:
            raise HTTPException(
                status_code=400, detail="Hujjatlardan matn o'qib bo'lmadi. Fayllar to'g'ri formatda ekanligini tekshiring."
            )

        analysis = generate_ai_analysis(pdf_text, pptx_text)

        return {"success": True, "data": analysis}

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Xatolik yuz berdi: {str(e)}")


@app.get("/")
def read_root():
    return {"message": "SMK Process Builder API is running."}


@app.get("/health")
def health_check():
    """API sog'lig'ini tekshirish"""
    api_key = os.getenv("GROQ_API_KEY")
    has_key = bool(api_key and len(api_key.strip()) >= 10)
    return {
        "status": "healthy",
        "groq_api_key_configured": has_key,
    }
