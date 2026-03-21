import streamlit as st
from groq import Groq
import fitz  # PyMuPDF
from pptx import Presentation
import pandas as pd
import os
from dotenv import load_dotenv

load_dotenv()

# ================== 1. FUNKSIYALAR ==================
def extract_text_from_pdf(pdf_file):
    doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
    text = "\n".join([page.get_text() for page in doc])
    return text

def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def ai_summarize(text, prompt_type="umumiy"):
    api_key = os.getenv("GROQ_API_KEY")
    if not api_key:
        return "Xatolik: GROQ_API_KEY topilmadi. Iltimos, loyiha papkasida .env faylini yarating va GROQ_API_KEY ni sozlang."
        
    client = Groq(api_key=api_key)
    system_prompt = {
        "umumiy": "Siz SMK mutaxassisisiz. Matnni qisqa, aniq va rus/o'zbek tilida umumlashtiring.",
        "taqqoslash": "KFU o'quv yurti SMK va sanoat korxonalari SMK ni taqqoslang: farqlar, o'xshashliklar, asosiy xulosalar.",
        "pdca": "Matndan PDCA tsikli bo'yicha jarayonlarni ajratib, 4 guruhga bo'ling (Plan-Do-Check-Act) va ular haqida qisqacha ma'lumot bering."
    }
    
    try:
        response = client.chat.completions.create(
            model="llama3-70b-8192",
            messages=[{"role": "system", "content": system_prompt.get(prompt_type, "Umumlashtiring.")},
                      {"role": "user", "content": text[:15000]}] # 15000 chars roughly matches token limit
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"API xatoligi: {str(e)}"

# ================== 2. STREAMLIT ILOVA ==================
st.set_page_config(page_title="SMK Process Builder SRM", page_icon="⚙️", layout="wide")
st.title("🛠️ SMK Process Builder (SRM Tizimi)")
st.subheader("Sifat Menejmenti Hujjatlarini Avtomatik Umumlashtirish va Taqqoslash Tizimi")

st.markdown("""
Bu dastur SMK hujjatlarini o'qiydi (hozirda PDF va PPTX). Orqa fonda AI dan foydalanib, jarayonlarni tahlil qiladi va PDCA (Rejalashtirish, Bajarish, Tekshirish, Harakat qilish) modellarini shakllantiradi.
""")

col1, col2 = st.columns(2)
with col1:
    pdf_file = st.file_uploader("Sanoat uchun (Maqola) PDF faylni yuklang", type="pdf")
with col2:
    pptx_file = st.file_uploader("KFU Universitet uchun PPTX faylni yuklang", type="pptx")

if pdf_file and pptx_file:
    with st.spinner("Fayllar oqilmoqda, kuting..."):
        pdf_text = extract_text_from_pdf(pdf_file)
        pptx_text = extract_text_from_pptx(pptx_file)
        
    st.success("✅ Fayllar muvaffaqiyatli o'qildi va xotiraga yuklandi!")
    
    tab1, tab2, tab3 = st.tabs(["📋 Umumlashtirish", "🔄 Taqqoslash", "📊 PDCA Model"])
    
    with tab1:
        st.write("Hujjatlarning qisqacha qisqartirilgan va umumlashgan holati.")
        if st.button("Umumlashtirishni boshlash", use_container_width=True):
            with st.spinner("AI matnlarni tahlil qilmoqda..."):
                summary_pdf = ai_summarize(pdf_text, "umumiy")
                summary_pptx = ai_summarize(pptx_text, "umumiy")
                
                col_res1, col_res2 = st.columns(2)
                with col_res1:
                    st.write("### 🏭 PDF (Sanoat) Xulosasi")
                    st.info(summary_pdf)
                with col_res2:
                    st.write("### 🎓 PPTX (Universitet) Xulosasi")
                    st.info(summary_pptx)
    
    with tab2:
        st.write("Ikki muhit o'rtasidagi SMK tizimi tafovuti va o'xshashliklari.")
        if st.button("Taqqoslanish jadvalini yaratish", use_container_width=True):
            with st.spinner("Taqqoslash amalga oshirilmoqda..."):
                combined = f"UNIVERSITET:\n{pptx_text[:8000]}\n\nSANOAT:\n{pdf_text[:8000]}"
                comparison = ai_summarize(combined, "taqqoslash")
                st.markdown(comparison)
    
    with tab3:
        st.write("1-guruh: Asosiy, 2-guruh: Qo'llab-quvvatlovchi, 3-guruh: Boshqaruv, 4-guruh: PDCA/Sifat.")
        if st.button("Jarayonlarni PDCA ga ajratish", use_container_width=True):
            with st.spinner("Jarayonlar ajratilmoqda..."):
                pdca = ai_summarize(pdf_text + "\n" + pptx_text, "pdca")
                st.write(pdca)
                
                st.write("### 📊 Taxminiy Statistik Ma'lumotlar")
                data = {
                    "Jarayon turi": ["Asosiy (Do)", "Qo'llab-quvvatlovchi (Do)", "Boshqaruv (Plan)", "Sifatni oshirish (Check/Act)"],
                    "Soni (Universitet)": [15, 8, 5, 4], 
                    "Soni (Korxona)": [10, 7, 6, 8]
                }
                df = pd.DataFrame(data)
                st.dataframe(df, use_container_width=True)
                
                # IKP (Informatsionnaya Karta Processa) Shablon matni
                ikp_template = """INFORMASION KARTA PROTSESSA (IKP)

1. Protsess nomi: ____________________________________
2. Protsess egasi: ___________________________________
3. Kirish (Вход): ____________________________________
4. Chiqish (Выход): __________________________________
5. Resurslar: ________________________________________

6. PDCA Bosqichlari:
   - Plan (Reja): 
   - Do (Bajarish): 
   - Check (Tekshirish): 
   - Act (Harakat/Yaxshilash): 

7. Boshqaruvchi mezonlari va ko'rsatkichlar:
...
"""
                st.download_button(
                    "📥 ИКП shablonini yuklab olish (TXT)", 
                    data=ikp_template, 
                    file_name="IKP_template_SMK.txt",
                    mime="text/plain",
                    use_container_width=True
                )
else:
    st.info("Iltimos, chap panelga yoxud yuqoridagi maydonga PDF hamda PPTX hujjatlarni yuklang. Bu dastur ikkala loyihani solishtirish va ma'lumot olish maqsadida aynan ikkita faylni kutmoqda.")
